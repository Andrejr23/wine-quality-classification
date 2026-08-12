# -*- coding: utf-8 -*-
"""Gera a apresentacao executiva do Tech Challenge Fase 2 - Wine Quality Classification."""

import math
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_TICK_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path(r"C:\Users\deehj\OneDrive\FIAP\wine-quality-classification\apresentacao\Apresentacao_Executiva_Wine_Quality.pptx")

# ---------------------------------------------------------------- design system
BERRY = RGBColor(0x6D, 0x2E, 0x46)
ROSE = RGBColor(0xA2, 0x67, 0x69)
CREAM = RGBColor(0xEC, 0xE2, 0xD0)
CREAM_SOFT = RGBColor(0xF7, 0xF1, 0xE8)
DARK = RGBColor(0x2B, 0x12, 0x19)
DARK_CARD = RGBColor(0x43, 0x21, 0x2A)
INK = RGBColor(0x3A, 0x2C, 0x30)
MUTED = RGBColor(0x8A, 0x7A, 0x7E)
MUTED_LIGHT = RGBColor(0xC9, 0xB3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Cambria"
BODY_FONT = "Calibri"

W, H = 13.333, 7.5
ML = 0.6
CONTENT_W = W - 2 * ML  # 12.133

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def new_slide(bg=WHITE, notes=""):
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def tb(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    """paras: lista de dicts {text, size, color, bold, italic, font, align, space_after, line}"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(paras, dict):
        paras = [paras]
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        if p.get("space_after") is not None:
            para.space_after = Pt(p["space_after"])
        if p.get("line"):
            para.line_spacing = p["line"]
        run = para.add_run()
        run.text = p["text"]
        f = run.font
        f.name = p.get("font", BODY_FONT)
        f.size = Pt(p.get("size", 14))
        f.bold = p.get("bold", False)
        f.italic = p.get("italic", False)
        f.color.rgb = p.get("color", INK)
    return box


def card(slide, x, y, w, h, fill=CREAM_SOFT, line=None, radius=0.06):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    sh.adjustments[0] = radius
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def panel(slide, x, y, w, h, fill=CREAM):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def badge(slide, x, y, d, label, fill=BERRY, color=WHITE, size=15):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = BODY_FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return sh


def deco_circles(slide, cx, cy, radii, colors, transparency=None):
    """Circulos concentricos decorativos (motivo visual do deck)."""
    for r, c in zip(radii, colors):
        sh = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r)
        )
        sh.fill.solid()
        sh.fill.fore_color.rgb = c
        sh.line.fill.background()
        sh.shadow.inherit = False


TITLE_CHARS_PER_LINE = 52  # Cambria 33pt bold na largura util de 12,13"


def slide_title(slide, text, kicker=None, color=BERRY, kicker_color=MUTED, y=0.55):
    # O kicker acompanha a altura real do titulo: titulo de duas linhas empurra o subtitulo.
    lines = max(1, math.ceil(len(text) / TITLE_CHARS_PER_LINE))
    tb(slide, ML, y, CONTENT_W, lines * 0.62,
       {"text": text, "size": 33, "bold": True, "color": color, "font": TITLE_FONT, "line": 1.0})
    if kicker:
        tb(slide, ML, y + lines * 0.60 + 0.32, CONTENT_W - 1.0, 0.5,
           {"text": kicker, "size": 15.5, "color": kicker_color})


def eyebrow(slide, x, y, w, text, color=RGBColor(0xC9, 0x8A, 0x92)):
    tb(slide, x, y, w, 0.3, {"text": text.upper(), "size": 11.5, "bold": True, "color": color})


def style_chart(chart, colors, label_color=INK, value_axis=True, data_labels=True,
                number_format='0"%"'):
    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 60
    for i, ser in enumerate(plot.series):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = colors[i % len(colors)]
        ser.format.line.fill.background()
    if data_labels:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(12)
        dl.font.bold = True
        dl.font.name = BODY_FONT
        dl.font.color.rgb = label_color
        dl.number_format = number_format
        dl.number_format_is_linked = False
    ca = chart.category_axis
    ca.has_major_gridlines = False
    ca.tick_labels.font.size = Pt(12)
    ca.tick_labels.font.name = BODY_FONT
    ca.tick_labels.font.color.rgb = label_color
    ca.format.line.color.rgb = RGBColor(0xDD, 0xD3, 0xD6)
    va = chart.value_axis
    va.visible = value_axis
    # Com o eixo de valores oculto, as gridlines viram ruido - os rotulos ja dao a leitura.
    va.has_major_gridlines = value_axis
    if value_axis:
        va.major_gridlines.format.line.color.rgb = RGBColor(0xEC, 0xE6, 0xE8)
        va.major_gridlines.format.line.width = Pt(0.75)
        va.tick_labels.font.size = Pt(11)
        va.tick_labels.font.color.rgb = MUTED
    return chart


def color_point(chart, series_idx, point_idx, color):
    """Destaca um ponto especifico da serie (python-pptx pinta a serie inteira por padrao)."""
    pt = chart.plots[0].series[series_idx].points[point_idx]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = color
    pt.format.line.fill.background()


PLACEHOLDER_NOTE = "Numeros ilustrativos. Substituir pelos resultados reais apos rodar os notebooks."


# =============================================================== SLIDE 1 - Capa
s = new_slide(DARK, notes=(
    "Abertura (0:00-0:30). Apresente o grupo e a pergunta central: da para prever a qualidade "
    "de um vinho antes de alguem prova-lo? Preencher os nomes dos 5 integrantes."))
deco_circles(s, 10.6, 3.9, [2.75, 1.95, 1.15], [DARK_CARD, RGBColor(0x5A, 0x28, 0x38), BERRY])
eyebrow(s, 0.85, 1.5, 6.0, "Tech Challenge · Fase 2 · Machine Learning Applied to Business")
tb(s, 0.85, 2.0, 7.9, 2.2, [
    {"text": "Classificando a qualidade de vinhos com Machine Learning",
     "size": 40, "bold": True, "color": WHITE, "font": TITLE_FONT, "line": 1.05},
])
tb(s, 0.85, 4.35, 7.9, 0.45,
   {"text": "Pós-Tech FIAP · Data Analytics · Turma 14DTAT", "size": 16, "color": CREAM})
tb(s, 0.85, 5.15, 7.9, 1.3, [
    {"text": "INTEGRANTES", "size": 11, "bold": True, "color": ROSE, "space_after": 6},
    {"text": "Nome 1  ·  Nome 2  ·  Nome 3  ·  Nome 4  ·  Nome 5",
     "size": 14, "color": MUTED_LIGHT},
])
tb(s, 0.85, 6.6, 7.9, 0.35,
   {"text": "Setembro de 2026", "size": 12, "color": MUTED_LIGHT})

# =============================================================== SLIDE 2 - O problema
s = new_slide(notes=(
    "Contexto (0:30-1:15). O metodo atual e humano: caro, lento e inconsistente. "
    "Nao ataque o especialista - o ponto e que ele nao escala."))
slide_title(s, "A qualidade do vinho depende de quem prova",
            "A avaliação é sensorial: um especialista prova cada amostra e atribui uma nota.")
tb(s, ML, 2.45, 5.4, 3.4, [
    {"text": "Esse é o padrão da indústria há décadas — e funciona. O problema não é a "
             "competência do enólogo: é que o método não acompanha a escala da produção.",
     "size": 15.5, "color": INK, "line": 1.3, "space_after": 14},
    {"text": "Cada lote avaliado consome tempo de um profissional escasso, e a nota final "
             "carrega a variação natural de qualquer julgamento humano.",
     "size": 15.5, "color": INK, "line": 1.3},
])
rows = [
    ("01", "Subjetivo", "A nota depende do paladar, da memória sensorial e da experiência de quem prova."),
    ("02", "Lento", "Cada amostra exige uma sessão de degustação dedicada."),
    ("03", "Inconsistente", "Dois especialistas podem discordar sobre o mesmo lote."),
]
for i, (num, head, desc) in enumerate(rows):
    y = 2.35 + i * 1.42
    badge(s, 6.5, y, 0.72, num, size=14)
    tb(s, 7.45, y + 0.02, 5.25, 0.35,
       {"text": head, "size": 17, "bold": True, "color": BERRY})
    tb(s, 7.45, y + 0.42, 5.25, 0.8,
       {"text": desc, "size": 14, "color": MUTED, "line": 1.25})

# =============================================================== SLIDE 3 - A oportunidade
s = new_slide(notes=(
    "A virada do argumento: o dado que resolveria o problema ja existe. Ninguem precisa "
    "instalar sensor novo - a producao ja mede tudo isso por controle de qualidade."))
panel(s, 0, 0, 5.3, H, CREAM)
eyebrow(s, 0.8, 2.15, 3.9, "Já medidos em cada lote", BERRY)
tb(s, 0.8, 2.5, 3.9, 1.5,
   {"text": "11", "size": 90, "bold": True, "color": BERRY, "font": TITLE_FONT, "line": 0.9})
tb(s, 0.8, 4.1, 3.9, 1.6,
   {"text": "indicadores físico-químicos que a vinícola já coleta por controle de qualidade, "
            "antes de qualquer degustação.", "size": 15.5, "color": INK, "line": 1.3})
tb(s, 6.0, 1.15, 6.73, 1.5,
   {"text": "A informação que falta já está na linha de produção",
    "size": 33, "bold": True, "color": BERRY, "font": TITLE_FONT, "line": 1.05})
tb(s, 6.0, 2.85, 6.73, 1.0,
   {"text": "Acidez, açúcar, álcool, densidade e enxofre são medidos de rotina. A proposta é "
            "usar esses números para antecipar a nota que o especialista daria.",
    "size": 15.5, "color": INK, "line": 1.3})
gains = [
    ("Antecipa", "O produtor sabe o resultado provável antes do engarrafamento."),
    ("Padroniza", "Critério único para todos os lotes, sem variação de avaliador."),
    ("Direciona", "Aponta qual variável ajustar no processo para melhorar o lote."),
]
for i, (head, desc) in enumerate(gains):
    y = 4.2 + i * 0.95
    badge(s, 6.0, y, 0.52, str(i + 1), size=13)
    tb(s, 6.7, y - 0.02, 6.03, 0.32, {"text": head, "size": 15.5, "bold": True, "color": BERRY})
    tb(s, 6.7, y + 0.32, 6.03, 0.5, {"text": desc, "size": 13.5, "color": MUTED, "line": 1.2})

# =============================================================== SLIDE 4 - Os dados
s = new_slide(notes=(
    "Apresente a base sem jargao. Nao diga 'dataset' nem 'variaveis': diga 'medicoes'. "
    "Preencher o numero de amostras com o valor real da base escolhida."))
slide_title(s, "O que foi analisado",
            "Amostras de vinho com as medições de produção e a nota final dada pelos especialistas.")
stats = [("1.599", "amostras analisadas"), ("11", "medições por amostra"), ("0 a 10", "escala da nota final")]
for i, (big, lab) in enumerate(stats):
    x = ML + i * 4.15
    tb(s, x, 2.05, 3.9, 0.75,
       {"text": big, "size": 40, "bold": True, "color": BERRY, "font": TITLE_FONT, "line": 1.0})
    tb(s, x, 2.82, 3.9, 0.35, {"text": lab, "size": 13, "color": MUTED})
groups = [
    ("Acidez", "Acidez fixa, acidez volátil e ácido cítrico — definem o frescor e denunciam defeitos."),
    ("Açúcar residual", "O açúcar que sobra após a fermentação; separa vinhos secos de suaves."),
    ("Cloretos e sulfatos", "Sais presentes no vinho; influenciam salinidade e estabilidade."),
    ("Dióxido de enxofre", "Livre e total — é o conservante que protege o vinho da oxidação."),
    ("Densidade e pH", "Indicadores físicos ligados ao teor de álcool e ao equilíbrio ácido."),
    ("Teor alcoólico", "Resultado direto da maturação da uva e do ponto de colheita."),
]
for i, (head, desc) in enumerate(groups):
    col, row = i % 3, i // 3
    x = ML + col * 4.05
    y = 3.55 + row * 1.75
    card(s, x, y, 3.7, 1.55)
    tb(s, x + 0.25, y + 0.2, 3.2, 0.32, {"text": head, "size": 15, "bold": True, "color": BERRY})
    tb(s, x + 0.25, y + 0.6, 3.2, 0.85, {"text": desc, "size": 12.5, "color": MUTED, "line": 1.2})

# =============================================================== SLIDE 5 - Achado 1
s = new_slide(notes=(
    "Achado 1 - o desbalanceamento. Este e o slide mais importante da analise. "
    "Explique que vinho excelente e raro, e que por isso 'acertar a maioria' nao significa nada. "
    "Substituir os percentuais pelos valores reais."))
slide_title(s, "Achado 1 — vinho excelente é exceção, não regra",
            "A distribuição das notas mostra uma concentração forte na faixa intermediária.")
chart_data = CategoryChartData()
chart_data.categories = ["Baixa / Média\n(nota abaixo de 7)", "Alta qualidade\n(nota 7 ou mais)"]
chart_data.add_series("Participação", (86.4, 13.6))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(2.25),
                        Inches(6.3), Inches(4.3), chart_data)
style_chart(gf.chart, [ROSE], number_format='0.0"%"', value_axis=False)
color_point(gf.chart, 0, 1, BERRY)  # destaca a classe minoritaria
tb(s, 0.55, 6.62, 6.3, 0.3,
   {"text": PLACEHOLDER_NOTE, "size": 10.5, "italic": True, "color": MUTED})
card(s, 7.25, 2.25, 5.48, 2.15, CREAM)
tb(s, 7.6, 2.5, 4.8, 1.0,
   {"text": "1 em cada 7", "size": 40, "bold": True, "color": BERRY, "font": TITLE_FONT, "line": 1.0})
tb(s, 7.6, 3.42, 4.8, 0.8,
   {"text": "é a proporção aproximada de vinhos classificados como de alta qualidade.",
    "size": 14, "color": INK, "line": 1.25})
tb(s, 7.25, 4.75, 5.48, 1.9, [
    {"text": "Por que isso muda tudo", "size": 17, "bold": True, "color": BERRY, "space_after": 8},
    {"text": "Um método que dissesse “este vinho é comum” para absolutamente todos os lotes "
             "acertaria a grande maioria dos casos — e seria completamente inútil.",
     "size": 14, "color": INK, "line": 1.3, "space_after": 8},
    {"text": "Por isso o critério de sucesso aqui não é “quantos acertos”, e sim quantos vinhos "
             "realmente bons o método consegue encontrar.",
     "size": 14, "color": INK, "line": 1.3},
])

# =============================================================== SLIDE 6 - Achado 2
s = new_slide(notes=(
    "Achado 2 - o que os vinhos bons tem em comum. Fale de alcool e de sulfatos em linguagem "
    "de producao (ponto de colheita, maturacao da uva). Substituir os valores."))
slide_title(s, "Achado 2 — o que os vinhos bons têm em comum",
            "Medições que crescem junto com a nota atribuída pelos especialistas.")
cd = CategoryChartData()
cd.categories = ["Teor alcoólico", "Sulfatos", "Ácido cítrico", "Acidez fixa"]
cd.add_series("Relação com a nota", (0.48, 0.25, 0.23, 0.12))
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.55), Inches(2.25),
                        Inches(6.5), Inches(4.0), cd)
style_chart(gf.chart, [BERRY], number_format="0.00", value_axis=False)
gf.chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
tb(s, 0.55, 6.4, 6.5, 0.5,
   {"text": "Quanto maior a barra, mais forte a associação com uma nota alta. "
            + PLACEHOLDER_NOTE, "size": 10.5, "italic": True, "color": MUTED, "line": 1.2})
tb(s, 7.45, 2.3, 5.28, 4.3, [
    {"text": "O álcool é o sinal mais forte", "size": 19, "bold": True, "color": BERRY,
     "space_after": 10},
    {"text": "Teor alcoólico mais alto costuma vir de uvas mais maduras e concentradas — "
             "e é justamente esse corpo que o especialista premia na degustação.",
     "size": 14.5, "color": INK, "line": 1.3, "space_after": 16},
    {"text": "Leitura para a produção", "size": 19, "bold": True, "color": BERRY,
     "space_after": 10},
    {"text": "O ponto de colheita da uva aparece como uma das alavancas mais relevantes de "
             "qualidade — uma decisão que acontece meses antes do engarrafamento.",
     "size": 14.5, "color": INK, "line": 1.3, "space_after": 16},
    {"text": "Os sulfatos também aparecem entre os sinais positivos, sugerindo que a dosagem "
             "correta do conservante acompanha lotes mais bem avaliados.",
     "size": 14.5, "color": INK, "line": 1.3},
])

# =============================================================== SLIDE 7 - Achado 3
s = new_slide(notes=(
    "Achado 3 - o defeito. Acidez volatil e o acido acetico: e literalmente o caminho do vinho "
    "para virar vinagre. Este e o achado mais acionavel para a producao."))
slide_title(s, "Achado 3 — o que derruba a nota é um defeito",
            "A acidez volátil se comporta como o oposto da qualidade: quanto maior, pior a avaliação.")
comp = [
    (BERRY, WHITE, CREAM, "Vinhos de alta qualidade",
     ["Acidez volátil baixa e controlada",
      "Teor alcoólico acima da média do lote",
      "Boa presença de ácido cítrico e sulfatos"]),
    (CREAM_SOFT, BERRY, MUTED, "Vinhos comuns",
     ["Acidez volátil elevada — risco de acetificação",
      "Teor alcoólico mais baixo",
      "Menor concentração aromática e estrutural"]),
]
for i, (bg, head_c, body_c, head, items) in enumerate(comp):
    x = ML + i * 6.23
    card(s, x, 2.35, 5.9, 3.4, bg)
    tb(s, x + 0.4, 2.65, 5.1, 0.45, {"text": head, "size": 19, "bold": True, "color": head_c})
    for j, it in enumerate(items):
        y = 3.3 + j * 0.78
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.4), Inches(y + 0.09),
                                 Inches(0.13), Inches(0.13))
        dot.fill.solid()
        dot.fill.fore_color.rgb = head_c
        dot.line.fill.background()
        dot.shadow.inherit = False
        tb(s, x + 0.72, y, 4.78, 0.68, {"text": it, "size": 14, "color": body_c, "line": 1.25})
tb(s, ML, 6.05, CONTENT_W, 0.85, [
    {"text": "O que isso significa na prática", "size": 15.5, "bold": True, "color": BERRY,
     "space_after": 5},
    {"text": "Acidez volátil alta é o ácido acético — o mesmo do vinagre. Não é uma questão de "
             "gosto pessoal: é um defeito de processo, e defeito de processo se corrige.",
     "size": 14, "color": INK, "line": 1.25},
])

# =============================================================== SLIDE 8 - Correlacoes
s = new_slide(notes=(
    "Este slide atende a exigencia do enunciado de justificar cada correlacao. "
    "No video, cite so duas; o resto fica na apresentacao escrita. Preencher os valores."))
slide_title(s, "As relações entre as medições fazem sentido físico",
            "Cada associação encontrada nos dados tem uma explicação conhecida da enologia.")
pairs = [
    ("Álcool e densidade", "0,00", "O álcool é menos denso que a água: quanto mais álcool, menor a densidade do vinho. É uma relação física direta."),
    ("Acidez fixa e pH", "0,00", "Mais ácido significa pH menor. As duas medidas descrevem o mesmo fenômeno por ângulos opostos."),
    ("Enxofre livre e total", "0,00", "O enxofre livre é parte do total, então as duas medidas caminham juntas por definição."),
    ("Acidez volátil e nota", "0,00", "Quanto mais ácido acético, pior a avaliação — o defeito é percebido imediatamente na degustação."),
]
for i, (head, val, desc) in enumerate(pairs):
    col, row = i % 2, i // 2
    x = ML + col * 6.23
    y = 2.35 + row * 2.1
    card(s, x, y, 5.9, 1.9)
    tb(s, x + 0.4, y + 0.25, 3.7, 0.35, {"text": head, "size": 16, "bold": True, "color": BERRY})
    tb(s, x + 4.2, y + 0.18, 1.35, 0.45,
       {"text": val, "size": 20, "bold": True, "color": ROSE, "font": TITLE_FONT,
        "align": PP_ALIGN.RIGHT})
    tb(s, x + 0.4, y + 0.72, 5.1, 0.95, {"text": desc, "size": 13.5, "color": MUTED, "line": 1.25})
tb(s, ML, 6.65, CONTENT_W, 0.35,
   {"text": "Preencher cada valor com o coeficiente calculado na análise. " + PLACEHOLDER_NOTE,
    "size": 10.5, "italic": True, "color": MUTED})

# =============================================================== SLIDE 9 - Como testamos
s = new_slide(notes=(
    "Como testamos, sem matematica. A frase-chave: separamos parte dos vinhos e escondemos "
    "a nota deles, para conferir se o metodo acertaria sozinho."))
slide_title(s, "Como testamos",
            "Nenhum vinho usado para ensinar o método foi usado para avaliá-lo.")
steps = [
    ("01", "Separamos", "Uma parte dos vinhos foi reservada e teve a nota escondida do método."),
    ("02", "Ensinamos", "Com o restante, o método aprendeu a associar medições a notas altas."),
    ("03", "Comparamos", "Testamos dois caminhos diferentes de análise e confrontamos os resultados."),
    ("04", "Conferimos", "Revelamos as notas escondidas e medimos os acertos."),
]
w_step = 2.808
for i, (num, head, desc) in enumerate(steps):
    x = ML + i * (w_step + 0.3)
    card(s, x, 2.5, w_step, 2.85)
    badge(s, x + 0.35, 2.8, 0.72, num, size=14)
    tb(s, x + 0.35, 3.78, w_step - 0.7, 0.35, {"text": head, "size": 17, "bold": True, "color": BERRY})
    tb(s, x + 0.35, 4.22, w_step - 0.7, 1.0, {"text": desc, "size": 13, "color": MUTED, "line": 1.2})
    if i < 3:
        tb(s, x + w_step - 0.02, 3.65, 0.34, 0.4,
           {"text": "→", "size": 18, "bold": True, "color": ROSE, "align": PP_ALIGN.CENTER})
tb(s, ML, 5.7, CONTENT_W, 1.0, [
    {"text": "Por que dois caminhos e não um só", "size": 16, "bold": True, "color": BERRY,
     "space_after": 6},
    {"text": "Um método simples e transparente, que mostra claramente o peso de cada medição, e "
             "um método mais sofisticado, capaz de captar combinações que o primeiro não enxerga. "
             "Comparar os dois evita confiar em um resultado sem ter parâmetro.",
     "size": 14, "color": INK, "line": 1.3},
])

# =============================================================== SLIDE 10 - Resultado
s = new_slide(DARK, notes=(
    "Resultado (2:45-3:45). Nao fale em acuracia, F1 ou AUC. Traduza: 'de cada 10 vinhos bons, "
    "o metodo encontra X'. Substituir os numeros pelos resultados reais."))
tb(s, ML, 0.7, CONTENT_W, 0.85,
   {"text": "O resultado", "size": 33, "bold": True, "color": WHITE, "font": TITLE_FONT, "line": 1.0})
tb(s, ML, 1.6, CONTENT_W - 1.0, 0.5,
   {"text": "Medimos o que realmente importa: quantos vinhos de qualidade o método encontra.",
    "size": 15.5, "color": MUTED_LIGHT})
metrics = [
    ("00%", "dos vinhos de alta qualidade\nforam identificados corretamente"),
    ("00%", "dos vinhos apontados como bons\nrealmente eram bons"),
    ("00", "medições bastam para chegar\na esse resultado"),
]
for i, (big, lab) in enumerate(metrics):
    x = ML + i * 4.21
    card(s, x, 2.6, 3.71, 2.2, DARK_CARD)
    tb(s, x + 0.35, 2.9, 3.0, 0.85,
       {"text": big, "size": 44, "bold": True, "color": CREAM, "font": TITLE_FONT, "line": 1.0})
    tb(s, x + 0.35, 3.85, 3.0, 0.75, {"text": lab, "size": 13, "color": MUTED_LIGHT, "line": 1.25})
tb(s, ML, 5.2, CONTENT_W, 1.4, [
    {"text": "Uma ressalva honesta", "size": 17, "bold": True, "color": ROSE, "space_after": 8},
    {"text": "O método erra — e o tipo de erro importa. Deixar de reconhecer um vinho excelente "
             "significa vender abaixo do valor. Apontar como excelente um vinho comum expõe a "
             "marca. A escolha entre proteger um risco ou outro é uma decisão de negócio, "
             "não uma decisão técnica.",
     "size": 14.5, "color": CREAM, "line": 1.35},
])
tb(s, ML, 6.85, CONTENT_W, 0.3,
   {"text": PLACEHOLDER_NOTE, "size": 10.5, "italic": True, "color": MUTED_LIGHT})

# =============================================================== SLIDE 11 - O que fazer
s = new_slide(notes=(
    "Recomendacao (3:45-4:40). Aqui e onde o trabalho vira valor. Cada linha e uma acao concreta "
    "para a producao, nao uma constatacao."))
slide_title(s, "O que a vinícola pode fazer com isso",
            "Cada medição influente aponta para uma decisão concreta na produção.")
actions = [
    ("Acidez volátil", "Derruba a qualidade",
     "Reforçar higiene e controle microbiológico na fermentação para evitar acetificação.", BERRY),
    ("Teor alcoólico", "Eleva a qualidade",
     "Monitorar maturação da uva e ajustar o ponto de colheita lote a lote.", ROSE),
    ("Dióxido de enxofre", "Precisa de equilíbrio",
     "Calibrar a dosagem: protege da oxidação, mas em excesso compromete o aroma.", RGBColor(0x8A, 0x4A, 0x5C)),
]
for i, (var, efeito, acao, color) in enumerate(actions):
    y = 2.45 + i * 1.5
    card(s, ML, y, CONTENT_W, 1.32)
    badge(s, ML + 0.35, y + 0.34, 0.64, str(i + 1), fill=color, size=15)
    tb(s, ML + 1.2, y + 0.28, 2.9, 0.35, {"text": var, "size": 17, "bold": True, "color": BERRY})
    tb(s, ML + 1.2, y + 0.68, 2.9, 0.32, {"text": efeito, "size": 13, "color": MUTED})
    tb(s, ML + 4.4, y + 0.42, 7.3, 0.7, {"text": acao, "size": 15, "color": INK, "line": 1.25})
tb(s, ML, 6.95, CONTENT_W, 0.35,
   {"text": "Ajustar os itens acima conforme os resultados reais da análise.",
    "size": 10.5, "italic": True, "color": MUTED})

# =============================================================== SLIDE 12 - Proximos passos
s = new_slide(DARK, notes=(
    "Fechamento (4:40-5:00). Encerre com a frase de impacto e seja honesto sobre os limites - "
    "isso aumenta a credibilidade diante da diretoria."))
deco_circles(s, 11.6, 6.4, [2.3, 1.55, 0.85], [DARK_CARD, RGBColor(0x5A, 0x28, 0x38), BERRY])
tb(s, ML, 0.8, 8.5, 1.3,
   {"text": "Próximos passos", "size": 33, "bold": True, "color": WHITE, "font": TITLE_FONT,
    "line": 1.0})
nexts = [
    ("Validar em produção", "Aplicar o método a lotes novos e comparar com a avaliação do especialista."),
    ("Ampliar a base", "Incluir safras, regiões e castas diferentes para testar a generalização."),
    ("Apoiar, não substituir", "O método prioriza o que merece atenção; a palavra final continua sendo do enólogo."),
]
for i, (head, desc) in enumerate(nexts):
    y = 2.35 + i * 1.25
    badge(s, ML, y, 0.62, str(i + 1), fill=BERRY, size=14)
    tb(s, ML + 0.95, y - 0.02, 7.4, 0.35, {"text": head, "size": 18, "bold": True, "color": CREAM})
    tb(s, ML + 0.95, y + 0.4, 7.4, 0.6, {"text": desc, "size": 14, "color": MUTED_LIGHT, "line": 1.25})
tb(s, ML, 6.15, 9.3, 0.9,
   {"text": "A degustação continua sendo arte. O que mudamos é o momento em que a vinícola "
            "descobre o resultado.",
    "size": 19, "italic": True, "color": WHITE, "font": TITLE_FONT, "line": 1.25})

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print("Salvo:", OUT)
print("Slides:", len(prs.slides.__iter__.__self__._sldIdLst))
